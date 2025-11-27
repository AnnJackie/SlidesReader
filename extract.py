import os
import re

from pptx import Presentation

def extract_text_no_duplicates(path):
    prs = Presentation(path)
    all_slides = []

    def extract_from_shape(shape, bucket):
        # Text frames (text boxes, titles, etc.)
        if shape.has_text_frame:
            text_no_tabs = shape.text_frame.text.strip().replace('\x0b', '').replace('\t', '')
            text_no_tabs_no_extra_newlines = re.sub(r'\n+', '\n', text_no_tabs)
            text_no_tabs_no_extra_newlines_no_extra_spaces = re.sub(r'\s+', ' ', text_no_tabs_no_extra_newlines)
            if text_no_tabs_no_extra_newlines_no_extra_spaces:
                bucket.append(text_no_tabs_no_extra_newlines_no_extra_spaces)

        # Tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        bucket.append(cell_text)

        # Grouped shapes
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for shp in shape.shapes:
                extract_from_shape(shp, bucket)

    for slide_number, slide in enumerate(prs.slides, start=1):
        bucket = []
        for shape in slide.shapes:
            extract_from_shape(shape, bucket)

        # Remove duplicates while preserving order
        seen = set()
        unique_text = []
        for t in bucket:
            if t not in seen and t != str(slide_number):
                seen.add(t)
                unique_text.append(t)

        all_slides.append({
            "slide_number": slide_number,
            "text": unique_text
        })

    return all_slides


presentations = {}
for ppt_name in os.listdir('slides'):
    presentation = extract_text_no_duplicates(os.path.join('slides', ppt_name))
    name = ppt_name.split('.')[0]
    presentations[name] = presentation

for name in presentations:
    with open(os.path.join('texts', f'{name}.txt'), 'w', encoding='utf-8') as txt_file:
        for s in presentations[name]:
            txt_file.write(f"{s['slide_number']}.\n")
            for t in s["text"]:
                txt_file.write(f'{t}\n')